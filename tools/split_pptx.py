#!/usr/bin/env python3
"""
PPTX 拆分工具：将 PPTX 中每个独立形状拆分为单独的 PPTX 文件。

用法:
    python split_pptx.py <source.pptx> [output_dir] [选项]

选项:
    --list               仅列出所有形状信息（不导出），用于排查
    --slides 1,3,5       只处理指定页码（逗号分隔，从 1 开始）
    --no-frame           跳过无填充的形状（如边框、方框）
    --min-width 100      跳过宽度小于 N（px）的形状
    --min-height 100     跳过高度小于 N（px）的形状
    --skip-text          跳过纯文本框形状

示例:
    python split_pptx.py res.pptx              # 导出所有形状
    python split_pptx.py res.pptx --list       # 先看看有哪些形状
    python split_pptx.py res.pptx out/ --no-frame  # 跳过方框
"""

import sys
import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

SHAPE_TYPE_NAMES = {
    MSO_SHAPE_TYPE.AUTO_SHAPE: 'auto_shape',
    MSO_SHAPE_TYPE.GROUP: 'group',
    MSO_SHAPE_TYPE.PICTURE: 'picture',
    MSO_SHAPE_TYPE.FREEFORM: 'freeform',
    MSO_SHAPE_TYPE.TEXT_BOX: 'text_box',
    MSO_SHAPE_TYPE.PLACEHOLDER: 'placeholder',
}

def describe_shape(shape):
    stype = SHAPE_TYPE_NAMES.get(shape.shape_type, f'other({shape.shape_type})')
    w = emu_to_px(shape.width)
    h = emu_to_px(shape.height)
    name = shape.name
    has_fill = True
    try:
        fill = shape.fill
        has_fill = fill.type is not None
    except Exception:
        pass
    text = shape.text[:30] if shape.has_text_frame and shape.text else ''
    return f"  [{stype}] {name} ({w}x{h}px) fill={'Y' if has_fill else 'N'} text={repr(text)}"

def emu_to_px(emu):
    return round(emu / 12700)

def is_text_only_shape(shape):
    """Check if shape is just a text box with no visual content."""
    if shape.has_text_frame:
        # Has text but check if it's a placeholder or has fill/line
        if not shape.has_text_frame:
            return False
        # If it's just a plain text box with text and no fill/line
        try:
            fill = shape.fill
            line = shape.line
            has_fill = fill.type is not None  # type is None = no fill
            has_line = line.fill.type is not None if line else False
            if not has_fill and not has_line:
                return True
        except Exception:
            pass
    return False

def list_shapes(source_path, slides=None):
    """Print all shapes without exporting."""
    prs = Presentation(source_path)
    target = list(range(len(prs.slides))) if slides is None else [s - 1 for s in slides]
    for sidx in target:
        slide = prs.slides[sidx]
        print(f"\n第 {sidx+1} 页 ({len(slide.shapes)} 个形状):")
        for i, s in enumerate(slide.shapes):
            print(f"  [{i+1}] {describe_shape(s)}")

def shape_has_no_fill(shape):
    try:
        return shape.fill.type is None
    except Exception:
        return False

def split_pptx(source_path, output_dir, slides=None, min_width=0, min_height=0,
               skip_text=False, list_only=False, no_frame=False):
    prs = Presentation(source_path)

    if list_only:
        list_shapes(source_path, slides)
        return

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    os.makedirs(output_dir, exist_ok=True)

    total_slides = len(prs.slides)
    target_slides = list(range(total_slides)) if slides is None else [s - 1 for s in slides]

    print(f"源文件: {source_path}")
    print(f"共 {total_slides} 页，处理第 {[s+1 for s in target_slides]} 页")

    total = 0
    skipped = 0

    for slide_idx in target_slides:
        slide = prs.slides[slide_idx]
        shapes = list(slide.shapes)
        print(f"\n第 {slide_idx+1} 页: {len(shapes)} 个形状")

        for shape_idx, shape in enumerate(shapes):
            # Apply filters
            w_px = emu_to_px(shape.width)
            h_px = emu_to_px(shape.height)

            if min_width and w_px < min_width:
                skipped += 1
                continue
            if min_height and h_px < min_height:
                skipped += 1
                continue
            if skip_text and is_text_only_shape(shape):
                skipped += 1
                continue
            if no_frame and shape_has_no_fill(shape):
                skipped += 1
                continue

            try:
                new_prs = Presentation()
                new_prs.slide_width = slide_width
                new_prs.slide_height = slide_height

                blank = new_prs.slide_layouts[6]
                new_slide = new_prs.slides.add_slide(blank)

                el = deepcopy(shape._element)
                new_slide.shapes._spTree.append(el)
                new_shape = new_slide.shapes[-1]

                new_shape.left = slide_width // 2 - shape.width // 2
                new_shape.top = slide_height // 2 - shape.height // 2

                output_name = f"slide{slide_idx+1:02d}_shape{shape_idx+1:03d}.pptx"
                new_prs.save(os.path.join(output_dir, output_name))

                total += 1
                if total % 20 == 0:
                    print(f"  已导出 {total} 个...")
            except Exception as e:
                print(f"  ✗ shape {shape_idx+1}: {e}")

    print(f"\n总计: 导出 {total} 个 | 跳过 {skipped} 个 → {output_dir}/")

def parse_args():
    args = {'slides': None, 'min_width': 0, 'min_height': 0, 'skip_text': False,
            'no_frame': False, 'list_only': False}
    i = 2
    while i < len(sys.argv):
        a = sys.argv[i]
        if a == '--slides' and i + 1 < len(sys.argv):
            args['slides'] = [int(s.strip()) for s in sys.argv[i + 1].split(',')]
            i += 2
        elif a == '--min-width' and i + 1 < len(sys.argv):
            args['min_width'] = int(sys.argv[i + 1])
            i += 2
        elif a == '--min-height' and i + 1 < len(sys.argv):
            args['min_height'] = int(sys.argv[i + 1])
            i += 2
        elif a == '--skip-text':
            args['skip_text'] = True
            i += 1
        elif a == '--no-frame':
            args['no_frame'] = True
            i += 1
        elif a == '--list':
            args['list_only'] = True
            i += 1
        else:
            i += 1
    return args

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    src = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 and not sys.argv[2].startswith('--') else 'split_output'
    opts = parse_args()

    split_pptx(src, out, **opts)
    print("\n下一步: 将这些 .pptx 文件拖入 PPTArts 导入页面批量导入")
